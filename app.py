import streamlit as st
import pandas as pd
import gspread
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import json
import io
import datetime
import copy
import time
import re
import ast # Required for robust parsing of single-quoted AI data

# --- PAGE SETUP ---
st.set_page_config(page_title="WBR Tool", page_icon="🚀", layout="wide")
st.title("🚀 Strategic Performance Dashboard")

# --- 🔒 1. PASSWORD PROTECTION (MOVED TO SIDEBAR) ---
if "APP_PASSWORD" not in st.secrets:
    st.sidebar.error("⚠️ Password not set in Secrets. Please add APP_PASSWORD to secrets.toml")
    st.stop()

password = st.sidebar.text_input("Enter Password", type="password")
if not password:
    st.sidebar.info("Please enter the password to proceed.")
    st.stop()
if password != st.secrets["APP_PASSWORD"]:
    st.sidebar.error("❌ Incorrect password.")
    st.stop()

# --- 🔗 2. URL INPUT (MOVED TO SIDEBAR) ---
sheet_url = st.sidebar.text_input("Google Sheet URL", placeholder="Paste Link Here...")
if not sheet_url:
    st.sidebar.warning("Please paste the Google Sheet URL to connect.")
    st.stop()

st.markdown("---")

# --- CONFIGURATION & THRESHOLDS ---
HISTORY_TAB_NAME = "WBR_History"
TEMPLATE_FILE = "template.pptx"
AHT_GOAL_SEC = 780  # 13 minutes standard

# --- AUTHENTICATION ---
@st.cache_resource
def connect_services(target_url):
    genai.configure(api_key=st.secrets["GEMINI_KEY"])
    try:
        models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
        model_name = next((m for m in models if 'flash' in m), models[0])
        ai = genai.GenerativeModel(model_name)
    except:
        ai = genai.GenerativeModel("gemini-1.5-flash")
    
    secrets_val = st.secrets["GAA_JSON"]
    if isinstance(secrets_val, str):
        creds_dict = json.loads(secrets_val)
    else:
        creds_dict = dict(secrets_val)

    if "private_key" in creds_dict:
        creds_dict["private_key"] = creds_dict["private_key"].replace("\\n", "\n")

    gc = gspread.service_account_from_dict(creds_dict)
    sh = gc.open_by_url(target_url)
    return ai, sh

try:
    model, sh = connect_services(sheet_url)
    st.success("✅ Connected")
except Exception as e:
    st.error(f"❌ Connection Failed: {e}")
    st.stop()

# --- HELPERS ---
def clean_num(x):
    try: return float(str(x).replace('%','').replace(',',''))
    except: return 0

def get_first_name(full_name):
    if "," in full_name: return full_name.split(",")[1].strip()
    return full_name.split(" ")[0]

def sec_to_min(seconds):
    try:
        m, s = int(seconds // 60), int(seconds % 60)
        return f"{m}:{s:02d}"
    except: return "0:00"

def get_trend_arrow(curr, prev):
    if curr - prev > 0.01: return "⬆️"
    if curr - prev < -0.01: return "⬇️"
    return "➡️"

def get_color_for_score(value, is_iqa=True):
    if is_iqa:
        if value < 0.70: return RGBColor(255, 0, 0)
        if value < 0.85: return RGBColor(255, 165, 0)
        return RGBColor(0, 128, 0)
    else:
        if value < 0.80: return RGBColor(255, 0, 0)
        if value < 0.90: return RGBColor(255, 165, 0)
        return RGBColor(0, 128, 0)

def clean_markdown(text):
    if not text: return ""
    if isinstance(text, list): text = "\n".join([str(i) for i in text])
    text = str(text).replace("**", "").replace("##", "").replace("* ", "• ")
    return text.strip()

# --- MEMORY FUNCTIONS ---
def get_agent_history(agent_name):
    try:
        ws = sh.worksheet(HISTORY_TAB_NAME)
        data = ws.get_all_values()
        for row in reversed(data):
            if len(row) > 3 and row[1] == agent_name:
                return {"date": row[0], "focus": row[2], "plan": row[3]}
    except: return None
    return None

def save_agent_history(agent_name, focus, plan):
    try:
        ws = sh.worksheet(HISTORY_TAB_NAME)
        today = datetime.date.today().strftime("%Y-%m-%d")
        ws.append_row([today, agent_name, str(focus), str(plan)])
    except Exception as e: 
        print(f"Error saving history: {e}") # Changed to print error instead of silent pass

# --- STEP 1: LOAD DATA ---
if st.button("🔄 Analyze Data"):
    with st.spinner("Analyzing Trends & Identifying Stars..."):
        try:
            worksheet = sh.worksheet("MASTER_DATA")
            rows = worksheet.get_all_values()
            df = pd.DataFrame(rows[1:], columns=rows[0])
            
            # --- DATA CLEANING ---
            df.columns = df.columns.str.strip()
            
            # Identify the Notes column (Prioritize 'AI_Notes_Context')
            notes_col = 'AI_Notes_Context' if 'AI_Notes_Context' in df.columns else 'Notes'
            
            df['IQA_Score'] = df['IQA_Score'].apply(clean_num)
            df['Show_Rate'] = df['Show_Rate'].apply(clean_num)
            df['AHT'] = df['AHT'].apply(clean_num)
            df['Agent_Name'] = df['Agent_Name'].astype(str).str.strip()
            
            if 'Team' in df.columns:
                df['Team'] = df['Team'].astype(str).str.strip()
            
            if df['IQA_Score'].mean() > 1.5: df['IQA_Score'] /= 100
            if df['Show_Rate'].mean() > 1.5: df['Show_Rate'] /= 100
            
            # Smart Sort
            def get_sort_key(week_str):
                try:
                    parts = str(week_str).split('-')
                    return int(parts[0]) * 100 + int(parts[1])
                except: return 0
            
            df['Sort_Key'] = df['Week_ID'].apply(get_sort_key)
            
            candidates = []
            top_performers = [] # CHANGED to array to collect all top scorers
            agents = df['Agent_Name'].unique()
            
            for agent in agents:
                adf = df[df['Agent_Name'] == agent].sort_values('Sort_Key')
                if len(adf) < 3: continue
                w3, w2, w1 = adf.iloc[-3], adf.iloc[-2], adf.iloc[-1]
                
                # CHANGED: Capture multi-week notes to establish behavior trends
                notes_history = f"Wk-2: {w3.get(notes_col, 'None')} | Wk-1: {w2.get(notes_col, 'None')} | Curr: {w1.get(notes_col, 'None')}"
                
                comp = (w1['IQA_Score'] * 0.7) + (w1['Show_Rate'] * 0.3)
                
                # CHANGED: Collect all agents meeting Star Criteria for dropdown selection
                if w1['IQA_Score'] >= 0.95:
                    top_performers.append({
                        'Full_Name': agent, 'First_Name': get_first_name(agent),
                        'Trend': [w3['IQA_Score'], w2['IQA_Score'], w1['IQA_Score']],
                        'Current_IQA': w1['IQA_Score'], 'Current_SR': w1['Show_Rate'],
                        'Current_AHT': w1['AHT'], 'Notes': notes_history,
                        'Comp_Score': comp
                    })

                # --- DIAGNOSTIC INSIGHT LOGIC ---
                issues = []
                if w1['IQA_Score'] < 0.85: 
                    gap = 0.85 - w1['IQA_Score']
                    issues.append(f"Quality Gap (-{gap:.0%})")
                elif w2['IQA_Score'] < 0.85 or w3['IQA_Score'] < 0.85:
                    issues.append("Inconsistent Quality (Trend)")
                elif (w2['IQA_Score'] - w1['IQA_Score']) > 0.08: 
                    issues.append("Quality Regression")
                
                if w1['Show_Rate'] < 0.85: 
                    issues.append("Reliability Risk")
                if w1['AHT'] > AHT_GOAL_SEC:
                    over_by = w1['AHT'] - AHT_GOAL_SEC
                    issues.append(f"AHT Over (+{sec_to_min(over_by)})")
                
                if issues:
                    candidates.append({
                        'Full_Name': agent, 'First_Name': get_first_name(agent),
                        'Reason': ", ".join(issues),
                        'Trend': [w3['IQA_Score'], w2['IQA_Score'], w1['IQA_Score']],
                        'Current_IQA': w1['IQA_Score'], 'Current_SR': w1['Show_Rate'],
                        'Current_AHT': w1['AHT'], 'Notes': notes_history # CHANGED to pass multi-week notes
                    })
            
            st.session_state['candidates'] = sorted(candidates, key=lambda x: x['Current_IQA'])
            st.session_state['top_performers'] = sorted(top_performers, key=lambda x: x['Comp_Score'], reverse=True)
            st.session_state['df_team'] = df
            st.rerun()
        except Exception as e:
            st.error(f"Error reading sheet: {e}")

# --- STEP 2: SELECT & GENERATE ---
if 'candidates' in st.session_state:
    st.write("### 📋 Proposed Watchlist")
    formatted_preview = pd.DataFrame(st.session_state['candidates']).copy()
    formatted_preview['Current IQA'] = formatted_preview['Current_IQA'].map('{:.0%}'.format)
    st.table(formatted_preview[['Full_Name', 'Reason', 'Current IQA']]) 
    
    all_agent_names = [c['Full_Name'] for c in st.session_state['candidates']]
    select_all = st.checkbox("Select all identified agents", value=False)
    
    if select_all:
        selection = st.multiselect("Select Agents to Include:", all_agent_names, default=all_agent_names)
    else:
        selection = st.multiselect("Select Agents to Include:", all_agent_names, default=all_agent_names[:5])
    
    final_watchlist = [c for c in st.session_state['candidates'] if c['Full_Name'] in selection]
    
    # CHANGED: Dropdown Selection for Top Performer
    st.write("---")
    if 'top_performers' in st.session_state and st.session_state['top_performers']:
        star_names = [p['Full_Name'] for p in st.session_state['top_performers']]
        selected_star = st.selectbox("🌟 Select the Top Performer of the Week:", star_names)
        # Assign the selected agent to st.session_state['star'] for PowerPoint generation
        st.session_state['star'] = next((p for p in st.session_state['top_performers'] if p['Full_Name'] == selected_star), None)
    else:
        st.session_state['star'] = None
        st.info("No agents met the 95%+ Star criteria this week.")
    
    if st.button("🎨 Generate Strategic PowerPoint"):
        with st.spinner("Writing History & Generating Slides..."):
            prs = Presentation(TEMPLATE_FILE)

            def replace_text_colored(slide, replacements):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            for run in p.runs:
                                for k, (val, color) in replacements.items():
                                    if k in run.text:
                                        clean_val = clean_markdown(val)
                                        run.text = run.text.replace(k, clean_val)
                                        if color: run.font.color.rgb = color
                                        if k == "{{TEAM_INSIGHT_SUMMARY}}": run.font.size = Pt(13)
            
            df = st.session_state['df_team']
            replace_text_colored(prs.slides[0], {"{{DATE_TODAY}}": (datetime.date.today().strftime("%B %d, %Y"), None)})
            
            # --- 1. IDENTIFY TEAMS & PREPARE SLIDES ---
            original_dash_slide = prs.slides[1]
            if 'Team' in df.columns:
                unique_teams = [t for t in df['Team'].unique() if t and str(t).lower() not in ['nan', 'none', '']]
            else:
                unique_teams = []
            
            dashboards_to_process = []
            dashboards_to_process.append({
                'name': 'OVERALL', 
                'df': df, 
                'slide': original_dash_slide
            })
            
            for team in unique_teams:
                team_df = df[df['Team'] == team]
                new_slide = prs.slides.add_slide(original_dash_slide.slide_layout)
                for shape in original_dash_slide.shapes:
                    new_el = copy.deepcopy(shape.element)
                    new_slide.shapes._spTree.append(new_el)
                
                dashboards_to_process.append({
                    'name': team.upper(),
                    'df': team_df,
                    'slide': new_slide
                })

            # --- 2. PROCESS EACH DASHBOARD ---
            for ctx in dashboards_to_process:
                current_df = ctx['df']
                current_slide = ctx['slide']
                context_name = ctx['name']
                
                if not current_df.empty:
                    latest_wk_id = current_df['Week_ID'].unique()[-1]
                    latest_wk_data = current_df[current_df['Week_ID'] == latest_wk_id]
                    t_iqa = latest_wk_data['IQA_Score'].mean()
                    t_sr = latest_wk_data['Show_Rate'].mean()
                    t_aht = latest_wk_data['AHT'].mean()
                    target_hit = (latest_wk_data['AHT'] <= AHT_GOAL_SEC).mean()
                else:
                    t_iqa, t_sr, t_aht, target_hit = 0, 0, 0, 0

                try: 
                    team_summ_prompt = (
                        f"{context_name} Performance Review: IQA {t_iqa:.1%}, Show Rate {t_sr:.1%}, AHT {sec_to_min(t_aht)}. "
                        f"Critically, {target_hit:.0%} of the group met the 13m AHT standard. "
                        f"Task: Analyze correlation between quality and efficiency for {context_name}. STRICT MAX 40 WORDS. No asterisks."
                    )
                    team_summ = model.generate_content(team_summ_prompt).text.strip()
                except: team_summ = "Stable performance; focus remains on increasing AHT goal achievement rates."

                title_text = "WEEKLY PERFORMANCE REVIEW"
                if context_name != "OVERALL":
                    title_text = f"WEEKLY PERFORMANCE REVIEW - {context_name}"

                replace_text_colored(current_slide, {
                    "{{VAL_1}}": (f"{t_iqa:.1%}", get_color_for_score(t_iqa, True)),
                    "{{VAL_2}}": (f"{t_sr:.1%}", get_color_for_score(t_sr, False)),
                    "{{VAL_3}}": (sec_to_min(t_aht), None),
                    "{{TEAM_INSIGHT_SUMMARY}}": (team_summ, None),
                    "{{DATE_TODAY}}": (datetime.date.today().strftime("%B %d, %Y"), None),
                    "WEEKLY PERFORMANCE REVIEW": (title_text, None)
                })

                for shape in current_slide.shapes:
                    if shape.has_text_frame and "{{CHART_PLACEHOLDER}}" in shape.text:
                        left, top, w, h = shape.left, shape.top, shape.width, shape.height
                        shape._element.getparent().remove(shape._element) 
                        weeks = current_df['Week_ID'].unique()[-8:]
                        cdata = CategoryChartData()
                        cdata.categories = weeks
                        cdata.add_series('IQA Avg', [current_df[current_df['Week_ID'] == wk]['IQA_Score'].mean() for wk in weeks])
                        cdata.add_series('Show Rate', [current_df[current_df['Week_ID'] == wk]['Show_Rate'].mean() for wk in weeks])
                        cdata.add_series('AHT Goal %', [(current_df[current_df['Week_ID'] == wk]['AHT'] <= AHT_GOAL_SEC).mean() for wk in weeks])
                        chart = current_slide.shapes.add_chart(XL_CHART_TYPE.LINE, left, top, w, h, cdata).chart
                        chart.has_legend, chart.legend.position, chart.legend.font.size = True, XL_LEGEND_POSITION.BOTTOM, Pt(10)
                        chart.value_axis.maximum_scale = 1.0

            # --- SLIDE 3 (WATCHLIST) ---
            watchlist_slide = None
            for s in prs.slides:
                if any(sh.has_table for sh in s.shapes):
                    watchlist_slide = s
                    break
            
            if watchlist_slide:
                table = next((s.table for s in watchlist_slide.shapes if s.has_table), None)
                if table:
                    for i, item in enumerate(final_watchlist):
                        if i+1 >= len(table.rows): break
                        row = table.rows[i+1]
                        metric_category = "General"
                        if "Quality" in item['Reason']: metric_category = "Quality"
                        elif "AHT" in item['Reason']: metric_category = "Efficiency"
                        elif "Reliability" in item['Reason']: metric_category = "Reliability"
                        row.cells[0].text, row.cells[1].text = item['Full_Name'], metric_category
                        row.cells[2].text, row.cells[3].text = f"{item['Trend'][0]:.0%}", f"{item['Current_IQA']:.0%} {get_trend_arrow(item['Current_IQA'], item['Trend'][1])}"
                        row.cells[4].text, row.cells[5].text, row.cells[6].text = f"{(item['Current_IQA'] - item['Trend'][0]):.1%}", "90%", item['Reason']

            # --- DEEP DIVES ---
            template_slide = None
            template_idx = -1
            for idx, s in enumerate(prs.slides):
                for sh in s.shapes:
                    if sh.has_text_frame and "{{AGENT_NAME}}" in sh.text:
                        template_slide = s
                        template_idx = idx
                        break
                if template_slide: break
            
            TMP_IDX = template_idx if template_slide else 3

            def get_ai_analysis(mode, item):
                history = get_agent_history(item['Full_Name'])
                h_ctx = f"PREV: {history['focus']} | Plan: {history['plan']}" if history else "Baseline."
                
                # --- UPDATED PROMPTS FOR BETTER TREND/BEHAVIORAL ACCURACY ---
                if mode == "COACH":
                    prompt = (f"Expert Performance Coach. AGENT: {item['First_Name']} "
                              f"Data: IQA {item['Current_IQA']:.1%}, AHT {sec_to_min(item['Current_AHT'])}. "
                              f"Score Trend: {item['Trend'][0]:.0%} -> {item['Trend'][1]:.0%} -> {item['Trend'][2]:.0%}. "
                              f"3-Week Context/Notes: {item['Notes']}. " 
                              f"Standard: 13m goal. {h_ctx}. Identify specific behavioral trends from notes and explicitly state if performance is improving/declining. JSON keys MUST be 'analysis' (Max 50 words) and 'plan' (3-4 bullets). No asterisks.")
                else:
                    prompt = (f"Celebrate Performer {item['First_Name']} (IQA {item['Current_IQA']:.1%}). "
                              f"3-Week Context/Notes: {item['Notes']}. " 
                              f"TASK: 1. Winning Analysis (Max 50 words detailing exactly WHAT behaviors/actions they are doing well based on the notes). 2. Growth Plan (3 bullets). JSON FORMAT ONLY: {{'analysis': '...', 'plan': '...'}}. No other keys.")
                try:
                    res = model.generate_content(prompt)
                    json_match = re.search(r'\{.*\}', res.text, re.DOTALL)
                    if not json_match: raise ValueError("No JSON block found")
                    txt = json_match.group()
                    try: raw_data = json.loads(txt)
                    except: raw_data = ast.literal_eval(txt)
                    ans = raw_data.get('analysis') or "Review diagnostic metrics."
                    pln = raw_data.get('plan') or "• Monitor weekly trends."
                    if mode == "COACH": save_agent_history(item['Full_Name'], ans, pln)
                    return {"analysis": ans, "plan": pln}
                except: return {"analysis": "Manual Coaching Review Pending.", "plan": "• Review weekly performance trends."}

            for item in final_watchlist:
                res = get_ai_analysis("COACH", item)
                sl = prs.slides.add_slide(prs.slide_layouts[6])
                if template_slide:
                    for s in template_slide.shapes:
                        sl.shapes._spTree.append(copy.deepcopy(s.element))
                replace_text_colored(sl, {"{{AGENT_NAME}}": (item['Full_Name'], None), "{{SCORE}}": (f"{item['Current_IQA']:.0%} {get_trend_arrow(item['Current_IQA'], item['Trend'][1])}", get_color_for_score(item['Current_IQA'])), "{{AI_ANALYSIS_TEXT}}": (res['analysis'], None), "{{AI_ACTION_PLAN}}": (res['plan'], None)})

            if st.session_state['star']:
                star = st.session_state['star']
                res = get_ai_analysis("STAR", star)
                sl = prs.slides.add_slide(prs.slide_layouts[6])
                if template_slide:
                    for sh in template_slide.shapes:
                        sl.shapes._spTree.append(copy.deepcopy(sh.element))
                replace_text_colored(sl, {"{{AGENT_NAME}}": (f"⭐ {star['Full_Name']} ⭐", RGBColor(0,128,0)), "{{SCORE}}": (f"{star['Current_IQA']:.0%}", RGBColor(0,128,0)), "{{AI_ANALYSIS_TEXT}}": (res['analysis'], None), "{{AI_ACTION_PLAN}}": (res['plan'], None)})

            if template_slide:
                xml_slides = prs.slides._sldIdLst
                xml_slides.remove(list(xml_slides)[template_idx])
            
            binary_output = io.BytesIO()
            prs.save(binary_output)
            binary_output.seek(0)
            st.success("✨ Strategic Report Ready.")
            st.download_button(label="📥 Download PowerPoint", data=binary_output, file_name="WBR_Report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
